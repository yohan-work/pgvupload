import os, io, hashlib, requests
from typing import List, Optional

from fastapi import FastAPI, UploadFile, Form, HTTPException
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv
from supabase import create_client, Client

from pypdf import PdfReader
import docx as docx_reader
from pptx import Presentation

load_dotenv()

APP = FastAPI(title="RAG Upload Admin (FastAPI)")

# ===== CORS 설정 (Nextron 앱 연동용) =====
APP.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Nextron 앱에서 접근 가능
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ===== 환경설정 =====
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_ANON_KEY") 
OLLAMA_BASE = os.getenv("OLLAMA_BASE", "http://127.0.0.1:11434")
EMBED_MODEL = os.getenv("EMBED_MODEL", "nomic-embed-text")
CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "800"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "150"))

# Supabase 클라이언트 초기화
supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

# ===== 유틸 =====
def extract_text_from_file(file_bytes: bytes, filename: str, mime: Optional[str]) -> str:
    name_lower = (filename or "").lower()
    mime = (mime or "").lower()
    
    # PDF 파일
    if name_lower.endswith(".pdf") or "pdf" in mime:
        return extract_pdf_text(file_bytes)
    
    # PowerPoint 파일
    if (name_lower.endswith((".ppt", ".pptx")) or 
        "presentation" in mime or 
        "powerpoint" in mime or
        "presentationml" in mime):
        return extract_pptx_text(file_bytes)
    
    # Word 파일
    if (name_lower.endswith(".docx") or 
        "wordprocessing" in mime or 
        ("word" in mime and "presentation" not in mime)):
        return extract_docx_text(file_bytes)
    
    # 일반 텍스트
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except:
        return ""

def extract_pdf_text(file_bytes: bytes) -> str:
    text = []
    with io.BytesIO(file_bytes) as buf:
        reader = PdfReader(buf)
        for page in reader.pages:
            t = page.extract_text() or ""
            text.append(t)
    return "\n".join(text).strip()

def extract_docx_text(file_bytes: bytes) -> str:
    with io.BytesIO(file_bytes) as buf:
        doc = docx_reader.Document(buf)
        return "\n".join(p.text for p in doc.paragraphs).strip()

def extract_pptx_text(file_bytes: bytes) -> str:
    """PowerPoint 파일에서 텍스트 추출"""
    text = []
    with io.BytesIO(file_bytes) as buf:
        prs = Presentation(buf)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = shape.text.strip()
                    if shape_text:
                        text.append(shape_text)
                # 표 안의 텍스트도 추출
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                text.append(cell_text)
    return "\n".join(text).strip()

def chunk_text(text: str, size: int, overlap: int) -> List[str]:
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    N = len(text)
    while start < N:
        end = min(start + size, N)
        chunks.append(text[start:end])
        if end == N:
            break
        start = max(0, end - overlap)
    return [c.strip() for c in chunks if c.strip()]

def sha256_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()

def ollama_embed(texts: List[str]) -> List[List[float]]:
    out = []
    for t in texts:
        r = requests.post(
            f"{OLLAMA_BASE}/api/embeddings",
            json={"model": EMBED_MODEL, "prompt": t},
            timeout=120
        )
        r.raise_for_status()
        out.append(r.json()["embedding"])
    return out

def upsert_to_db(title: str, filename: str, mime: Optional[str], bytes_len: int,
                 whole_text: str, chunks: List[str], embeds: List[List[float]]) -> dict:
    doc_hash = sha256_text(whole_text)
    
    # 1) documents 테이블에 upsert
    doc_result = supabase.table("documents").upsert({
        "title": title,
        "source_path": filename,
        "mime_type": mime,
        "bytes": bytes_len,
        "hash": doc_hash
    }, on_conflict="hash").execute()
    
    if not doc_result.data:
        raise Exception("문서 저장 실패")
    
    document_id = doc_result.data[0]["id"]
    
    # 2) 기존 chunks 삭제
    supabase.table("chunks").delete().eq("document_id", document_id).execute()
    
    # 3) 새 chunks 삽입 (RPC 함수로 vector 변환)
    chunks_data = []
    for idx, (content, embedding) in enumerate(zip(chunks, embeds)):
        chunks_data.append({
            "document_id": document_id,
            "chunk_index": idx,
            "content": content,
            "embedding": embedding  # List[float] 그대로 전달
        })
    
    # RPC 함수로 embedding을 vector로 변환해서 삽입
    try:
        result = supabase.rpc("insert_chunks_with_embeddings", {
            "chunks_data": chunks_data
        }).execute()
        
        return {"document_id": str(document_id), "chunks": len(chunks)}
        
    except Exception as e:
        # RPC 함수가 없는 경우, 임시로 embedding을 JSON으로 저장
        print(f"RPC 함수 사용 실패, JSON으로 저장: {e}")
        
        # embedding을 JSON 문자열로 변환해서 저장 (임시)
        for chunk_data in chunks_data:
            chunk_data["embedding"] = str(chunk_data["embedding"])
        
        chunk_result = supabase.table("chunks").insert(chunks_data).execute()
        return {"document_id": str(document_id), "chunks": len(chunks)}

# ===== 뷰 (1페이지) =====
FORM_HTML = """
<!doctype html>
<html lang="ko">
<head>
  <meta charset="utf-8" />
  <title>RAG 업로드 어드민</title>
  <meta name="viewport" content="width=device-width,initial-scale=1" />
  <style>
    body { font-family: ui-sans-serif, system-ui, -apple-system, Segoe UI, Roboto, Arial; max-width: 720px; margin: 40px auto; padding: 0 16px; }
    .card { border: 1px solid #e5e7eb; border-radius: 12px; padding: 24px; }
    h1 { font-size: 20px; margin: 0 0 12px; }
    label { display:block; margin:12px 0 6px; font-weight:600; }
    input[type="text"], input[type="file"] { width:100%; padding:10px; border:1px solid #d1d5db; border-radius:8px; }
    button { margin-top:16px; padding:10px 16px; border-radius:10px; border:1px solid #111827; background:#111827; color:white; cursor:pointer; }
    .msg { margin-top:16px; white-space:pre-wrap; font-family:ui-monospace,Menlo,Consolas,monospace; }
    .hint { color:#6b7280; font-size:13px; }
  </style>
</head>
<body>
  <div class="card">
    <h1>RAG 업로드 어드민</h1>
    <form action="/upload" method="post" enctype="multipart/form-data">
      <label>문서 제목</label>
      <input type="text" name="title" placeholder="(선택) 미입력 시 파일명 사용" />
      <label>파일 선택</label>
      <input type="file" name="file" required />
      <div class="hint">PDF / DOCX / PPTX / TXT 지원</div>
      <button type="submit">업로드 → 인덱싱 → DB 저장</button>
    </form>
    <div class="msg">업로드 후 결과가 이 페이지로 표시됩니다.</div>
  </div>
</body>
</html>
"""

# ===== API 엔드포인트 (Nextron 앱용) =====

@APP.get("/api/status")
def api_status():
    """서비스 상태 확인"""
    try:
        # Ollama 연결 테스트
        ollama_status = requests.get(f"{OLLAMA_BASE}/api/tags", timeout=5).status_code == 200
        
        # Supabase 연결 테스트  
        supabase_status = bool(supabase.table("documents").select("id").limit(1).execute())
        
        return JSONResponse({
            "status": "healthy",
            "ollama": "connected" if ollama_status else "disconnected",
            "supabase": "connected" if supabase_status else "disconnected",
            "embed_model": EMBED_MODEL,
            "chunk_size": CHUNK_SIZE
        })
    except Exception as e:
        return JSONResponse({
            "status": "error",
            "error": str(e)
        }, status_code=500)

@APP.post("/api/upload")
async def api_upload(file: UploadFile, title: Optional[str] = Form(default=None)):
    """Nextron 앱에서 호출하는 JSON API"""
    try:
        file_bytes = await file.read()
        mime = file.content_type or "application/octet-stream"
        name = file.filename or "uploaded_file"
        title = title.strip() if title else name

        # 1) 텍스트 추출
        text = extract_text_from_file(file_bytes, name, mime)
        if not text.strip():
            raise HTTPException(400, f"텍스트 추출 실패: {name} ({mime})")

        # 2) 청크
        chunks = chunk_text(text, CHUNK_SIZE, CHUNK_OVERLAP)
        if not chunks:
            raise HTTPException(400, "청크 생성 실패(빈 텍스트)")

        # 3) 임베딩 (Ollama)
        embeds = ollama_embed(chunks)

        # 4) DB 저장
        result = upsert_to_db(
            title=title,
            filename=name,
            mime=mime,
            bytes_len=len(file_bytes),
            whole_text=text,
            chunks=chunks,
            embeds=embeds
        )
        
        return JSONResponse({
            "success": True,
            "document_id": result["document_id"],
            "title": title,
            "filename": name,
            "mime": mime,
            "bytes": len(file_bytes),
            "chunks": result["chunks"],
            "message": "문서 인덱싱 완료"
        })
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"처리 오류: {str(e)}")

@APP.get("/api/documents")
def api_list_documents(limit: int = 50, offset: int = 0):
    """업로드된 문서 목록"""
    try:
        result = supabase.table("documents") \
            .select("id, title, source_path, mime_type, bytes, created_at") \
            .order("created_at", desc=True) \
            .limit(limit) \
            .offset(offset) \
            .execute()
            
        return JSONResponse({
            "success": True,
            "documents": result.data,
            "count": len(result.data)
        })
    except Exception as e:
        raise HTTPException(500, f"문서 목록 조회 오류: {str(e)}")

@APP.post("/api/search-vectors")
async def api_search_vectors(request: dict):
    """Nextron용 벡터 데이터 제공 API"""
    try:
        query = request.get("query", "")
        
        if not query.strip():
            raise HTTPException(400, "검색어를 입력해주세요")
        
        # 1) 질의 임베딩 생성
        query_embedding = ollama_embed([query])[0]  # List[float]
        
        # 2) 모든 청크와 문서 정보 가져오기
        import json
        
        chunks_data = supabase.table("chunks") \
            .select("id, document_id, chunk_index, content, embedding") \
            .execute()
        
        docs_data = supabase.table("documents") \
            .select("id, title") \
            .execute()
        
        # 문서 title 매핑
        doc_titles = {doc["id"]: doc["title"] for doc in docs_data.data}
        
        # 청크 데이터 처리 (임베딩 파싱만)
        processed_chunks = []
        for chunk in chunks_data.data:
            try:
                # 문자열로 저장된 임베딩을 리스트로 변환
                if isinstance(chunk["embedding"], str):
                    embedding_str = chunk["embedding"].strip()
                    if embedding_str.startswith('[') and embedding_str.endswith(']'):
                        chunk_embedding = json.loads(embedding_str)
                    else:
                        chunk_embedding = [float(x.strip()) for x in embedding_str.split(',')]
                else:
                    chunk_embedding = chunk["embedding"]
                
                # 차원 검증
                if len(chunk_embedding) != len(query_embedding):
                    print(f"⚠️ 차원 불일치: query({len(query_embedding)}) vs chunk({len(chunk_embedding)})")
                    continue
                
                processed_chunks.append({
                    "id": chunk["id"],
                    "content": chunk["content"],
                    "embedding": chunk_embedding,  # 파싱된 벡터
                    "document_id": chunk["document_id"],
                    "document_title": doc_titles.get(chunk["document_id"], ""),
                    "chunk_index": chunk["chunk_index"]
                })
                    
            except Exception as parse_error:
                print(f"청크 {chunk['id']} 파싱 실패: {parse_error}")
                continue
        
        print(f"🔍 벡터 데이터 제공: '{query}' | 청크수={len(processed_chunks)} | 쿼리벡터차원={len(query_embedding)}")
        
        return JSONResponse({
            "success": True,
            "query": query,
            "query_embedding": query_embedding,  # 질의 벡터
            "chunks": processed_chunks,  # 모든 청크 데이터와 벡터
            "count": len(processed_chunks)
        })
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"벡터 데이터 조회 오류: {str(e)}")

@APP.post("/api/search")  
async def api_search(request: dict):
    """벡터 유사도 검색 (기존 방식 + 벡터 전용 옵션)"""
    try:
        query = request.get("query", "")
        top_k = request.get("top_k", 5)
        threshold = request.get("threshold", 0.5)
        vectors_only = request.get("vectors_only", False)  # 새 옵션
        
        if not query.strip():
            raise HTTPException(400, "검색어를 입력해주세요")
        
        # vectors_only 모드면 새로운 API로 리다이렉트
        if vectors_only:
            return await api_search_vectors(request)
        
        # 1) 질의 임베딩 생성
        query_embedding = ollama_embed([query])[0]  # List[float]
        
        # 2) Python으로 벡터 유사도 검색
        import json
        import math
        
        # 모든 청크와 문서 정보 가져오기
        chunks_data = supabase.table("chunks") \
            .select("id, document_id, chunk_index, content, embedding") \
            .execute()
        
        docs_data = supabase.table("documents") \
            .select("id, title") \
            .execute()
        
        # 문서 title 매핑
        doc_titles = {doc["id"]: doc["title"] for doc in docs_data.data}
        
        # 유사도 계산
        similarities = []
        for chunk in chunks_data.data:
            try:
                # 문자열로 저장된 임베딩을 리스트로 변환
                if isinstance(chunk["embedding"], str):
                    # "[1.0, 2.0, ...]" 형태의 문자열을 파싱
                    embedding_str = chunk["embedding"].strip()
                    if embedding_str.startswith('[') and embedding_str.endswith(']'):
                        chunk_embedding = json.loads(embedding_str)
                    else:
                        # "1.0,2.0,..." 형태면 split 사용
                        chunk_embedding = [float(x.strip()) for x in embedding_str.split(',')]
                else:
                    chunk_embedding = chunk["embedding"]
                
                # 차원 검증
                if len(chunk_embedding) != len(query_embedding):
                    print(f"⚠️ 차원 불일치: query({len(query_embedding)}) vs chunk({len(chunk_embedding)})")
                    continue
                
                # 코사인 유사도 계산
                dot_product = sum(a * b for a, b in zip(query_embedding, chunk_embedding))
                norm_query = math.sqrt(sum(a * a for a in query_embedding))
                norm_chunk = math.sqrt(sum(a * a for a in chunk_embedding))
                similarity = dot_product / (norm_query * norm_chunk) if (norm_query * norm_chunk) > 0 else 0
                
                # 모든 결과를 추가
                similarities.append({
                    "content": chunk["content"],
                    "similarity": float(similarity),
                    "document_id": chunk["document_id"],
                    "document_title": doc_titles.get(chunk["document_id"], ""),
                    "chunk_index": chunk["chunk_index"]
                })
                    
            except Exception as calc_error:
                print(f"청크 {chunk['id']} 유사도 계산 실패: {calc_error}")
                continue
        
        # 유사도 높은 순으로 정렬
        similarities.sort(key=lambda x: x["similarity"], reverse=True)
        
        # 임계값 필터링 및 상위 k개 선택
        results = [s for s in similarities if s["similarity"] > threshold][:top_k]
        
        # 디버그 정보 출력
        print(f"🔍 검색: '{query}' | 전체={len(similarities)} | 임계값>{threshold} | 결과={len(results)}")
        if similarities:
            print(f"📊 유사도: {similarities[0]['similarity']:.3f}(최고) ~ {similarities[-1]['similarity']:.3f}(최저)")
            for i, s in enumerate(similarities[:3]):
                print(f"  {i+1}. {s['similarity']:.3f}: {s['content'][:50]}...")
        
        return JSONResponse({
            "success": True,
            "query": query,
            "results": results,
            "count": len(results)
        })
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"검색 오류: {str(e)}")

# ===== 웹 페이지 (기존 유지) =====

@APP.get("/", response_class=HTMLResponse)
def index():
    return HTMLResponse(FORM_HTML)

@APP.post("/upload", response_class=HTMLResponse)
async def upload(file: UploadFile, title: Optional[str] = Form(default=None)):
    file_bytes = await file.read()
    mime = file.content_type or "application/octet-stream"
    name = file.filename or "uploaded_file"
    title = title.strip() if title else name

    # 1) 텍스트 추출
    text = extract_text_from_file(file_bytes, name, mime)
    if not text.strip():
        return HTMLResponse(f"<pre>텍스트 추출 실패: {name} ({mime})</pre>", status_code=400)

    # 2) 청크
    chunks = chunk_text(text, CHUNK_SIZE, CHUNK_OVERLAP)
    if not chunks:
        return HTMLResponse(f"<pre>청크 생성 실패(빈 텍스트)</pre>", status_code=400)

    # 3) 임베딩 (Ollama)
    try:
        embeds = ollama_embed(chunks)
    except Exception as e:
        return HTMLResponse(f"<pre>임베딩 오류: {str(e)}</pre>", status_code=500)

    # 4) DB 저장
    try:
        res = upsert_to_db(
            title=title,
            filename=name,
            mime=mime,
            bytes_len=len(file_bytes),
            whole_text=text,
            chunks=chunks,
            embeds=embeds
        )
    except Exception as e:
        return HTMLResponse(f"<pre>DB 저장 오류: {str(e)}</pre>", status_code=500)

    pretty = f"""
업로드/인덱싱 완료

- title: {title}
- filename: {name}
- mime: {mime}
- bytes: {len(file_bytes)}
- chunks: {res.get('chunks')}

document_id: {res.get('document_id')}
"""
    back = '<p><a href="/">← 돌아가기</a></p>'
    return HTMLResponse(f"<pre>{pretty}</pre>{back}")
